"""
PowerPoint File Handler
Translates .pptx files while preserving:
- Slide layouts and masters
- Text formatting (font, size, color, bold, italic)
- Shape positioning
- Tables
- Speaker notes
- Images (not modified)
"""

from pptx import Presentation
from core.utils import is_translatable


class PptxHandler:
    """Handles translation of .pptx files."""

    def translate(self, input_path, output_path, translator, progress_callback=None):
        """Translate all text content in a PowerPoint file."""
        prs = Presentation(input_path)

        total_slides = len(prs.slides)
        if total_slides == 0:
            prs.save(output_path)
            return

        for slide_idx, slide in enumerate(prs.slides):
            if progress_callback:
                progress = slide_idx / total_slides
                progress_callback(
                    progress,
                    f"Translating slide {slide_idx + 1} of {total_slides}..."
                )

            # Translate shapes
            for shape in slide.shapes:
                self._translate_shape(shape, translator)

            # Translate notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    if is_translatable(para.text):
                        translated = translator.translate_text(para.text)
                        self._set_paragraph_text(para, translated)

        prs.save(output_path)

        if progress_callback:
            progress_callback(1.0, "PowerPoint translation complete!")

    def _translate_shape(self, shape, translator):
        """Translate text in a shape."""
        # Regular text frames
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if is_translatable(para.text):
                    translated = translator.translate_text(para.text)
                    self._set_paragraph_text(para, translated)

        # Tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        if is_translatable(para.text):
                            translated = translator.translate_text(para.text)
                            self._set_paragraph_text(para, translated)

        # Group shapes (recursively)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child_shape in shape.shapes:
                    self._translate_shape(child_shape, translator)
            except Exception:
                pass

    def _set_paragraph_text(self, para, translated_text):
        """
        Set paragraph text while preserving first run's formatting.
        """
        if not para.runs:
            para.text = translated_text
            return

        first_run = para.runs[0]
        first_run.text = translated_text

        for run in para.runs[1:]:
            run.text = ""